import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # completely blank layout
    
    # ── Color Palette Definitions ──
    C_BG = RGBColor(10, 15, 30)          # Slate 955 / Dark Navy
    C_CARD = RGBColor(22, 29, 48)        # Slate 900 / Deep Slate
    C_BORDER = RGBColor(79, 70, 229)     # Indigo 600
    C_TITLE = RGBColor(248, 250, 252)    # White / Slate 50
    C_SUBTITLE = RGBColor(148, 163, 184) # Muted Slate Grey
    C_CYAN = RGBColor(6, 182, 212)       # Cyan Accent
    C_GREEN = RGBColor(16, 185, 129)     # Emerald Accent
    C_ORANGE = RGBColor(245, 158, 11)    # Amber Accent
    C_CODE_BG = RGBColor(15, 23, 42)     # Deep dark Slate for code blocks
    
    # Helper to set background
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = C_BG
        
    # Helper to draw accent top bar and slide header
    def add_header(slide, title, slide_num):
        # Top thin accent border
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = C_BORDER
        accent_bar.line.fill.background()
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = C_TITLE
        
        # Slide Number at bottom right
        num_box = slide.shapes.add_textbox(Inches(11.5), Inches(6.9), Inches(1.2), Inches(0.4))
        p_num = num_box.text_frame.paragraphs[0]
        p_num.alignment = PP_ALIGN.RIGHT
        p_num.text = f"Slide {slide_num} of 9"
        p_num.font.name = "Arial"
        p_num.font.size = Pt(9)
        p_num.font.color.rgb = C_SUBTITLE

    # Helper to parse and apply runs (supports **bold** to trigger Cyan highlight)
    def add_formatted_paragraph(tf, text, font_size=12, level=0, default_color=C_TITLE, bullet=True):
        p = tf.add_paragraph()
        p.level = level
        p.space_after = Pt(4)
        p.space_before = Pt(2)
        
        prefix = "•  " if bullet else ""
        text_to_parse = prefix + text
        
        parts = text_to_parse.split("**")
        for idx, part in enumerate(parts):
            if not part and idx == 0:
                continue
            run = p.add_run()
            run.text = part
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            if idx % 2 == 1:
                run.font.bold = True
                run.font.color.rgb = C_CYAN
            else:
                run.font.bold = False
                run.font.color.rgb = default_color
        return p

    # Helper to draw a standard container card
    def add_card(slide, title, left, top, width, height, border_color=C_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        
        if title:
            title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = "Arial"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = C_CYAN
            
        return card

    # Helper to write structured JSON fields inside a code box
    def add_code_block(slide, title, code_lines, left, top, width, height):
        # Card Background
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = C_CODE_BG
        card.line.color.rgb = C_BORDER
        card.line.width = Pt(1.5)
        
        # Card title
        title_box = slide.shapes.add_textbox(left, top - Inches(0.4), width, Inches(0.4))
        p_t = title_box.text_frame.paragraphs[0]
        p_t.text = f"📄 {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = C_CYAN
        
        # Code content
        code_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
        tf = code_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        
        for idx, line in enumerate(code_lines):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.space_after = Pt(1)
            p.space_before = Pt(0)
            
            # Simple color mapping for keys and values
            parts = line.split(":")
            if len(parts) >= 2 and not line.strip().startswith("//"):
                key_part = parts[0]
                value_part = ":".join(parts[1:])
                
                # Key run
                r1 = p.add_run()
                r1.text = key_part + ":"
                r1.font.name = "Consolas"
                r1.font.size = Pt(9.5)
                r1.font.color.rgb = RGBColor(148, 163, 184) # Slate gray keys
                
                # Value run
                r2 = p.add_run()
                r2.text = value_part
                r2.font.name = "Consolas"
                r2.font.size = Pt(9.5)
                if '"' in value_part:
                    r2.font.color.rgb = RGBColor(110, 231, 183) # Light green for strings
                elif "ObjectId" in value_part or "ISODate" in value_part:
                    r2.font.color.rgb = RGBColor(167, 139, 250) # Light purple for database functions
                else:
                    r2.font.color.rgb = RGBColor(253, 186, 116) # Orange for numbers/booleans
            else:
                # Normal or comment lines
                r = p.add_run()
                r.text = line
                r.font.name = "Consolas"
                r.font.size = Pt(9.5)
                if line.strip().startswith("//"):
                    r.font.color.rgb = RGBColor(100, 116, 139) # Dim gray for comments
                else:
                    r.font.color.rgb = RGBColor(248, 250, 252) # White for brackets

    # Helper to add standard grid tables
    def add_slide_table(slide, headers, rows, col_widths, left, top, width, height):
        table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
        table = table_shape.table
        
        # Apply column widths
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = w
            
        # Headers styling
        for idx, h in enumerate(headers):
            cell = table.cell(0, idx)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_BORDER
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(10.5)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
        # Rows styling
        for r_idx, row in enumerate(rows):
            for c_idx, text in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = text
                cell.fill.solid()
                if r_idx % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(18, 25, 41)
                else:
                    cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
                    
                p = cell.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(9.5)
                p.font.color.rgb = C_TITLE
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = C_CYAN

    # ==========================================
    # SLIDE 1: Title Slide (Modern Dark Theme)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)
    
    # Left vertical color bar
    color_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.6))
    color_bar.fill.solid()
    color_bar.fill.fore_color.rgb = C_BORDER
    color_bar.line.fill.background()
    
    # Title Text Frame
    title_box = s1.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "RareSense.AI"
    p1.font.name = "Arial"
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = C_TITLE
    p1.space_after = Pt(8)
    
    p2 = tf1.add_paragraph()
    p2.text = "LLM-Powered Rare Disease Detection from Unstructured Clinical Notes"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.color.rgb = C_CYAN
    p2.space_after = Pt(4)
    
    p3 = tf1.add_paragraph()
    p3.text = "Project Review II — Database Schema Design & Frontend Integration"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = C_SUBTITLE
    
    # Metadata Box
    meta_box = s1.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(8), Inches(2.2))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    
    add_formatted_paragraph(tf_meta, "**Candidate**: Jairaj Chilukala (B.Tech CSE, Woxsen University)", font_size=11, bullet=False)
    add_formatted_paragraph(tf_meta, "**Supervisor**: Raj Kumar Sir", font_size=11, bullet=False)
    add_formatted_paragraph(tf_meta, "**Domain / Focus**: Health Data Science (HDS) Systems", font_size=11, bullet=False)
    add_formatted_paragraph(tf_meta, "**Syllabus Equivalent**: Project #7 (CareerPulse Equivalent)", font_size=11, bullet=False)
    add_formatted_paragraph(tf_meta, "**Date**: June 2026", font_size=11, bullet=False)

    # ==========================================
    # SLIDE 2: [NEW POSITION] Recap of Review I
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "Context Recap: Review I Foundations", 2)
    
    # Left Card: Problem Statement
    add_card(s2, "The Diagnostic Odyssey", Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.2))
    prob_box = s2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_prob = prob_box.text_frame
    tf_prob.word_wrap = True
    add_formatted_paragraph(tf_prob, "**Diagnostic Delay**: Rare diseases take **7.6 years** and an average of **7 specialist visits** to diagnose because symptoms are scattered across separate systems.")
    add_formatted_paragraph(tf_prob, "**Context Isolation**: Specialists address isolated complaints (e.g. dermatologists treat a rash; hematologists treat low platelets) without connecting them chronologically.")
    add_formatted_paragraph(tf_prob, "**Review 1 Core Deliverable**: Concept mapping of a federated pipeline to unify patient timelines and run similarity matching.")
    
    # Right Card: Data Foundations
    add_card(s2, "Data Sources & Case Study", Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.2))
    ds_box = s2.shapes.add_textbox(Inches(7.05), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_ds = ds_box.text_frame
    tf_ds.word_wrap = True
    add_formatted_paragraph(tf_ds, "**EHR Data (MIMIC-IV)**: Leverages 40k+ real ICU records including discharge summaries, ICD codes, prescriptions, and lab events.")
    add_formatted_paragraph(tf_ds, "**Disease Ontology (HPO & Orphanet)**: Maps 7,000+ rare diseases to standardized phenotype symptom identifiers (e.g., malar rash = HP:0025300).")
    add_formatted_paragraph(tf_ds, "**Lupus Case Study Validation**: Showcased how uniting clinical notes from 4 separate doctor visits over 2 years flags Systemic Lupus Erythematosus (Lupus) in seconds.")

    # ==========================================
    # SLIDE 3: [NEW POSITION] User Interface & Dual-Modal Ingestion
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "Frontend Interface & Dual-Modal Ingestion Dashboard", 3)
    
    # Left Card: Dual-Modal Input Options
    add_card(s3, "Clinical Ingestion Panel (Text or Image)", Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.2))
    ingest_box = s3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_ingest = ingest_box.text_frame
    tf_ingest.word_wrap = True
    add_formatted_paragraph(tf_ingest, "**Free-Text Clinical Notes**: Text editor allows clinicians to copy-paste or write discharge reports, consultations, or progress updates.")
    add_formatted_paragraph(tf_ingest, "**Prescription Image Uploads**: Clinicians can upload **high-resolution images / scans** of handwritten prescriptions or lab printouts directly into the dashboard.")
    add_formatted_paragraph(tf_ingest, "**Vision-OCR Pipeline**: The system simulates an Optical Character Recognition (OCR) parser that processes images, extracts text boundaries, and pipes the raw text into the medSpaCy NLP extractor.")
    add_formatted_paragraph(tf_ingest, "**Instant Highlight Verification**: Extracted medical entities are highlighted on screen immediately, allowing clinician corrections before database saving.")

    # Right Card: UI Features & Dashboard Modules
    add_card(s3, "Clinician Dashboard Overview", Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.2))
    dash_box = s3.shapes.add_textbox(Inches(7.05), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_dash = dash_box.text_frame
    tf_dash.word_wrap = True
    add_formatted_paragraph(tf_dash, "**Patient Cohort Switcher**: Dropdown and filter tags to change patient profiles and admission IDs dynamically.")
    add_formatted_paragraph(tf_dash, "**Chronological Health Graph**: Graphic rendering of symptoms, drug timelines, and lab spikes over multi-year ranges.")
    add_formatted_paragraph(tf_dash, "**Recommendation list**: Ranked list of potential rare diseases showing similarity percentages (e.g. 87% match) with color-coded meters.")
    add_formatted_paragraph(tf_dash, "**Clinician Feedback Board**: Diagnostic validation checkboxes ('Approve' / 'Reject') and review text inputs to save clinician audits.")

    # ==========================================
    # SLIDE 4: [NEW POSITION] Database Design
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Database Schema: Document Architecture", 4)
    
    # Subtitle or explanation
    desc_box = s4.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.6))
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    p_desc = desc_tf.paragraphs[0]
    p_desc.text = "RareSense.AI structures its storage layer with MongoDB to process both textual clinical notes and uploaded prescription image metadata."
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = C_SUBTITLE
    
    # Grid Table detailing collections
    headers = ["Collection", "Primary Key / Indexes", "Document Purpose & Dual-Modal Ingestion Model", "Data Flow Relation"]
    rows = [
        ["patients", "_id (ObjectId)\nsubject_id (Unique)", "Stores base patient profiles, gender, age, and baseline demographics.", "Identity anchor"],
        ["clinical_notes", "_id (ObjectId)\nsubject_id (Index)\nnote_type (Index)", "Stores raw physician notes OR prescription image metadata (including image bucket paths and OCR-extracted texts).", "Unified Dual-Modal Ingestion feed"],
        ["patient_timelines", "_id (ObjectId)\nsubject_id (Unique)\nevents.timestamp (Index)", "Structured patient health history. Aggregates parsed events (symptoms, drugs, labs) and stores BioBERT vector embeddings.", "Aggregated timeline states"],
        ["rare_diseases", "_id (ObjectId)\norpha_code (Unique)", "Reference static dataset holding 7k+ disease-symptom maps and disease vectors.", "Orphanet reference matrix"],
        ["diagnosis_matches", "_id (ObjectId)\nsubject_id (Index)", "Historical prediction log storing similarity indexes and clinician review actions.", "Clinician validation audit loop"]
    ]
    col_widths = [Inches(1.8), Inches(2.2), Inches(5.2), Inches(2.9)]
    
    add_slide_table(s4, headers, rows, col_widths, Inches(0.6), Inches(1.8), Inches(12.133), Inches(4.8))

    # ==========================================
    # SLIDE 5: Schema JSON — Patients & Clinical Notes (Dual-Modal Ingest)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Schema Implementation: Patients & Ingested Notes/Images", 5)
    
    # Left Code Box: patients collection
    patients_json = [
        "{",
        "  \"_id\": ObjectId(\"666fca1234bc56de7890ef01\"),",
        "  \"subject_id\": 100234,",
        "  \"gender\": \"F\",",
        "  \"anchor_age\": 28,",
        "  \"anchor_year\": 2023,",
        "  \"created_at\": ISODate(\"2026-06-15T09:00:00Z\")",
        "}"
    ]
    add_code_block(s5, "patients Collection Schema", patients_json, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    
    # Right Code Box: clinical_notes collection (showing both image upload metadata and text notes)
    notes_json = [
        "{",
        "  \"_id\": ObjectId(\"666fca9876de54cb3210ab12\"),",
        "  \"subject_id\": 100234,",
        "  \"note_type\": \"prescription_image\", // OR \"text\"",
        "  \"image_url\": \"s3://hds-records/rx_100234_visit1.png\",",
        "  \"ocr_text\": \"28F malar rash... photosensitivity...\",",
        "  \"ocr_confidence\": 0.92,",
        "  \"ingested_date\": ISODate(\"2023-01-15T14:30:00Z\"),",
        "  \"metadata\": {",
        "    \"uploader_id\": \"DR_LIN\",",
        "    \"device\": \"Mobile Scan OCR\"",
        "  }",
        "}"
    ]
    add_code_block(s5, "clinical_notes Collection Schema (Dual-Modal)", notes_json, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 6: Schema JSON — Patient Timelines
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Schema Implementation: Aggregated Patient Timelines", 6)
    
    # Left Code Box: patient_timelines
    timeline_json = [
        "{",
        "  \"_id\": ObjectId(\"666fcafe54bc98de7610ef34\"),",
        "  \"subject_id\": 100234,",
        "  \"events\": [",
        "    {",
        "      \"event_id\": ObjectId(\"666fcb0123bc45de6789ff01\"),",
        "      \"timestamp\": ISODate(\"2023-01-15T14:30:00Z\"),",
        "      \"event_type\": \"symptom\",",
        "      \"source_note_id\": \"100234-NOTE-1\",",
        "      \"extracted_text\": \"malar rash\",",
        "      \"standardized_name\": \"Malar rash\",",
        "      \"hpo_code\": \"HP:0025300\",",
        "      \"confidence_score\": 0.94,",
        "      \"status\": \"present\"",
        "    },",
        "    {",
        "      \"event_id\": ObjectId(\"666fcb0123bc45de6789ff02\"),",
        "      \"timestamp\": ISODate(\"2023-08-22T10:15:00Z\"),",
        "      \"event_type\": \"lab\",",
        "      \"source_note_id\": \"100234-NOTE-2\",",
        "      \"extracted_text\": \"low platelets\",",
        "      \"standardized_name\": \"Thrombocytopenia\",",
        "      \"hpo_code\": \"HP:0001873\",",
        "      \"confidence_score\": 0.89,",
        "      \"status\": \"present\",",
        "      \"attributes\": { \"value\": 89000, \"unit\": \"uL\" }",
        "    }",
        "  ],",
        "  \"phenotype_vector\": [0.0125, -0.0432, 0.1194, 0.0053, \"...\"]",
        "}"
    ]
    add_code_block(s6, "patient_timelines Collection Schema", timeline_json, Inches(0.6), Inches(1.8), Inches(6.2), Inches(4.8))
    
    # Right Text Info Box
    add_card(s6, "Aggregate Timeline Logic", Inches(7.2), Inches(1.8), Inches(5.5), Inches(4.8))
    info_box = s6.shapes.add_textbox(Inches(7.35), Inches(2.3), Inches(5.2), Inches(4.1))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    
    add_formatted_paragraph(tf_info, "**Aggregate Timeline Model**: Compiles events from both raw text notes AND OCR-processed prescription text inputs into one consolidated history index.")
    add_formatted_paragraph(tf_info, "**Standardized Ontologies**: Resolves terms to Human Phenotype Ontology (HPO) and RxNorm codes to keep medical descriptions identical.")
    add_formatted_paragraph(tf_info, "**Negated Symptoms**: Parses exclusion terms (e.g. 'no evidence of rash') and stores them as negated status to prevent matching errors.")
    add_formatted_paragraph(tf_info, "**BioBERT Clinical Embeddings**: Computes 768-dimension vectors of aggregated symptoms for Atlas Vector Search.")

    # ==========================================
    # SLIDE 7: Schema JSON — Reference & Matches
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Schema Implementation: Reference Map & Matches", 7)
    
    # Left Code Box: rare_diseases
    diseases_json = [
        "{",
        "  \"_id\": ObjectId(\"666fcbde23ab45cd8912ef56\"),",
        "  \"orpha_code\": 536,",
        "  \"disease_name\": \"Systemic Lupus Erythematosus\",",
        "  \"synonyms\": [\"SLE\", \"Lupus\"],",
        "  \"symptom_profiles\": [",
        "    {",
        "      \"hpo_code\": \"HP:0025300\",",
        "      \"hpo_name\": \"Malar rash\",",
        "      \"frequency\": \"Very frequent (80-99%)\"",
        "    },",
        "    {",
        "      \"hpo_code\": \"HP:0001873\",",
        "      \"hpo_name\": \"Thrombocytopenia\",",
        "      \"frequency\": \"Frequent (30-79%)\"",
        "    }",
        "  ],",
        "  \"disease_vector\": [0.0142, -0.0398, 0.1215, 0.0049, \"...\"]",
        "}"
    ]
    add_code_block(s7, "rare_diseases Collection Schema", diseases_json, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    
    # Right Code Box: diagnosis_matches (with review status)
    matches_json = [
        "{",
        "  \"_id\": ObjectId(\"666fccaf12bc34de5678ef78\"),",
        "  \"subject_id\": 100234,",
        "  \"generated_at\": ISODate(\"2026-06-15T09:30:00Z\"),",
        "  \"matched_diseases\": [",
        "    {",
        "      \"rank\": 1,",
        "      \"orpha_code\": 536,",
        "      \"disease_name\": \"Systemic Lupus Erythematosus\",",
        "      \"matched_symptoms\": [\"HP:0025300\", \"HP:0001873\"],",
        "      \"vector_similarity_score\": 0.87,",
        "      \"rules_confidence_score\": 0.82,",
        "      \"final_confidence\": 0.85",
        "    }",
        "  ],",
        "  \"review_status\": \"approved\",",
        "  \"feedback_notes\": \"Highly consistent clinical timeline.\"",
        "}"
    ]
    add_code_block(s7, "diagnosis_matches Collection Schema", matches_json, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 8: Mock Pipeline Strategy (Backend-Less Parity)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Mock Integration Strategy & Clientside Prototyping", 8)
    
    # Left Card: Client-Side Simulation Pipeline
    add_card(s8, "Clientside Prototyping Architecture", Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.2))
    mock_box = s8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_mock = mock_box.text_frame
    tf_mock.word_wrap = True
    add_formatted_paragraph(tf_mock, "**Stored Local Datasets**: Embedded datasets representing clinical notes and Orphanet disease-symptom maps compiled into static Javascript lists.")
    add_formatted_paragraph(tf_mock, "**Mock OCR & Note Parser**: OCR and note parsing are simulated client-side via exact string indexing and concept dictionary lookups.")
    add_formatted_paragraph(tf_mock, "**Local Image Processing**: Drag-and-drop uploads of prescription images display an instant simulated OCR progress bar and dump text from pre-defined mock maps.")
    add_formatted_paragraph(tf_mock, "**localStorage Cache Hook**: Keeps track of clinician audits, notes, and verification statuses, simulating real CRUD updates.")

    # Right Card: Advantages of Client-side Simulation
    add_card(s8, "Integration Advantages", Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.2))
    adv_box = s8.shapes.add_textbox(Inches(7.05), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_adv = adv_box.text_frame
    tf_adv.word_wrap = True
    add_formatted_paragraph(tf_adv, "**Decoupled Development**: Allowed us to build, refine, and verify UI components and timeline visuals rapidly without database cluster reliance.")
    add_formatted_paragraph(tf_adv, "**1-to-1 Logical Equivalence**: The dashboard layout performs exactly like a server-connected app, ensuring a robust demonstration template.")
    add_formatted_paragraph(tf_adv, "**Zero Backend Runtime Needs**: Bypasses server configuration challenges for live evaluation, allowing immediate sandbox review.")

    # ==========================================
    # SLIDE 9: Review III Roadmap
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "Review III Production Roadmap", 9)
    
    # Step 1 Card
    add_card(s9, "1. DB & FastAPI Server Setup", Inches(0.6), Inches(1.8), Inches(3.8), Inches(4.3))
    s1_box = s9.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(3.6), Inches(3.6))
    tf_s1 = s1_box.text_frame
    tf_s1.word_wrap = True
    add_formatted_paragraph(tf_s1, "Set up active **MongoDB collections** for patient ingestion.")
    add_formatted_paragraph(tf_s1, "Build a FastAPI backend to handle REST API calls for note creation and uploads.")
    add_formatted_paragraph(tf_s1, "Connect image upload endpoints with storage nodes.")
    
    # Step 2 Card
    add_card(s9, "2. OCR & BioBERT Pipeline Integration", Inches(4.766), Inches(1.8), Inches(3.8), Inches(4.3))
    s2_box = s9.shapes.add_textbox(Inches(4.866), Inches(2.3), Inches(3.6), Inches(3.6))
    tf_s2 = s2_box.text_frame
    tf_s2.word_wrap = True
    add_formatted_paragraph(tf_s2, "Integrate **Tesseract OCR** or **Telescope-Vision LLM** for real prescription image scanning.")
    add_formatted_paragraph(tf_s2, "Generate real 768d vector embeddings from parsed timeline records using BioBERT.")
    add_formatted_paragraph(tf_s2, "Deploy NLP models inside a dedicated Docker container.")
    
    # Step 3 Card
    add_card(s9, "3. Live Vector Search Matching", Inches(8.933), Inches(1.8), Inches(3.8), Inches(4.3))
    s3_box = s9.shapes.add_textbox(Inches(9.033), Inches(2.3), Inches(3.6), Inches(3.6))
    tf_s3 = s3_box.text_frame
    tf_s3.word_wrap = True
    add_formatted_paragraph(tf_s3, "Enable native MongoDB Atlas **$vectorSearch** indexes.")
    add_formatted_paragraph(tf_s3, "Perform real-time cosine similarity matches on the database layer.")
    add_formatted_paragraph(tf_s3, "Validate accuracy metrics (sensitivity/specificity) against MIMIC-IV profiles.")
    
    # Export presentation
    filename = "RareSense_AI_Review2.pptx"
    try:
        prs.save(filename)
        print(f"[OK] PowerPoint presentation '{filename}' updated and generated successfully.")
    except PermissionError:
        fallback = "RareSense_AI_Review2_new.pptx"
        prs.save(fallback)
        print(f"[WARNING] '{filename}' is locked. Saved as '{fallback}' instead.")

if __name__ == "__main__":
    create_presentation()
